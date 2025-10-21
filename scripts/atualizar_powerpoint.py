"""
Atualiza PowerPoint - Adiciona Nomes dos Integrantes

Este script adiciona um slide com os nomes dos integrantes do grupo.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


def atualizar_powerpoint():
    """Adiciona slide com integrantes ao PowerPoint existente"""
    
    print("\n" + "=" * 70)
    print("  ATUALIZANDO POWERPOINT")
    print("=" * 70 + "\n")
    
    # Carregar apresentação existente
    print("📂 Carregando apresentação...")
    prs = Presentation('Apresentacao_Analise_Gastos_Parlamentares.pptx')
    
    # Adicionar novo slide após o título (posição 1)
    print("📊 Adicionando slide de integrantes...")
    slide_layout = prs.slide_layouts[6]  # Layout em branco
    slide = prs.slides.add_slide(slide_layout)
    
    # Mover slide para segunda posição
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[-1])
    xml_slides.insert(1, slides[-1])
    
    # Título
    left = Inches(1)
    top = Inches(1)
    width = Inches(8)
    height = Inches(0.8)
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "Integrantes do Grupo"
    
    p = tf.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.alignment = PP_ALIGN.CENTER
    
    # Integrantes
    integrantes = [
        'Leticia Cristina Silva - 21352',
        'Gabriel Davi Lopes Jacobini - 24734',
        'Thales Vinicius Leal Barcelos - 24740',
        'Maria Fernanda Leite Felicíssimo - 24767'
    ]
    
    left = Inches(2)
    top = Inches(2.5)
    width = Inches(6)
    height = Inches(4)
    
    txBox2 = slide.shapes.add_textbox(left, top, width, height)
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    
    for i, integrante in enumerate(integrantes):
        if i > 0:
            p = tf2.add_paragraph()
        else:
            p = tf2.paragraphs[0]
        
        p.text = integrante
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(0, 51, 102)
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(30)
    
    # Adicionar disciplina no rodapé
    left = Inches(1)
    top = Inches(6.5)
    width = Inches(8)
    height = Inches(0.5)
    
    txBox3 = slide.shapes.add_textbox(left, top, width, height)
    tf3 = txBox3.text_frame
    tf3.text = "Ciência de Dados - 2025"
    
    p3 = tf3.paragraphs[0]
    p3.font.size = Pt(14)
    p3.font.italic = True
    p3.font.color.rgb = RGBColor(100, 100, 100)
    p3.alignment = PP_ALIGN.CENTER
    
    # Salvar apresentação atualizada
    output_path = 'Apresentacao_Analise_Gastos_Parlamentares.pptx'
    prs.save(output_path)
    
    print(f"✅ PowerPoint atualizado com sucesso!")
    print(f"📄 Arquivo: {output_path}")
    print(f"📊 Total de slides: {len(prs.slides)}")
    
    return output_path


if __name__ == '__main__':
    try:
        arquivo = atualizar_powerpoint()
        
        print("\n" + "=" * 70)
        print("  🎉 SUCESSO!")
        print("=" * 70)
        print(f"\n  PowerPoint atualizado com slide de integrantes!")
        print(f"  Arquivo: {arquivo}\n")
        
    except Exception as e:
        print(f"\n❌ ERRO: {e}")
        import traceback
        traceback.print_exc()
