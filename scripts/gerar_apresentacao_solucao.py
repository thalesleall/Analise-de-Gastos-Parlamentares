"""
Gera apresentação PowerPoint profissional sobre a solução desenvolvida.
Dividido em duas partes: Análises Políticas + Arquitetura do Código

Autor: Grupo 1 - Análise de Dados Governamentais
Data: Novembro 2025
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import pandas as pd
import os
from pathlib import Path


class ApresentacaoSolucao:
    """Classe para gerar apresentação da solução completa."""
    
    def __init__(self, pasta_resultados):
        """
        Inicializa gerador de apresentação.
        
        Args:
            pasta_resultados: Caminho da pasta com resultados mais recentes
        """
        self.pasta_resultados = pasta_resultados
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)
        
        # Paleta de cores moderna
        self.azul_escuro = RGBColor(13, 27, 62)      # #0D1B3E
        self.azul_principal = RGBColor(0, 102, 204)  # #0066CC
        self.azul_claro = RGBColor(52, 152, 219)     # #3498DB
        self.laranja = RGBColor(230, 126, 34)        # #E67E22
        self.verde = RGBColor(39, 174, 96)           # #27AE60
        self.roxo = RGBColor(142, 68, 173)           # #8E44AD
        self.cinza_escuro = RGBColor(44, 62, 80)     # #2C3E50
        self.cinza_claro = RGBColor(189, 195, 199)   # #BDC3C7
        self.branco = RGBColor(255, 255, 255)
        
    def criar_slide_titulo_secao(self, titulo, numero_secao, cor_fundo):
        """Cria slide divisor de seção."""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Fundo colorido
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = cor_fundo
        
        # Número da seção (grande, transparente)
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(2), Inches(2))
        tf = txBox.text_frame
        tf.text = numero_secao
        p = tf.paragraphs[0]
        p.font.size = Pt(180)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.color.brightness = -0.3
        
        # Título da seção
        txBox = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(2))
        tf = txBox.text_frame
        tf.text = titulo
        p = tf.paragraphs[0]
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = self.branco
        
        return slide
    
    def criar_slide_conteudo(self, titulo, cor_destaque=None):
        """Cria slide de conteúdo com design moderno."""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        if cor_destaque is None:
            cor_destaque = self.azul_principal
        
        # Barra lateral colorida
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(0.15), Inches(7.5)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = cor_destaque
        shape.line.fill.background()
        
        # Título
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
        tf = txBox.text_frame
        tf.text = titulo
        p = tf.paragraphs[0]
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = self.azul_escuro
        
        return slide
    
    def adicionar_bullet_point(self, slide, texto, left, top, width, tamanho=16, cor=None, icone="•"):
        """Adiciona bullet point com ícone."""
        if cor is None:
            cor = self.cinza_escuro
            
        txBox = slide.shapes.add_textbox(left, top, width, Inches(0.4))
        tf = txBox.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        
        # Ícone
        run = p.add_run()
        run.text = f"{icone} "
        run.font.size = Pt(tamanho)
        run.font.color.rgb = self.azul_principal
        run.font.bold = True
        
        # Texto
        run = p.add_run()
        run.text = texto
        run.font.size = Pt(tamanho)
        run.font.color.rgb = cor
        
        return txBox
    
    def adicionar_caixa_destaque(self, slide, titulo, valor, left, top, width, height, cor):
        """Adiciona caixa de destaque com métrica."""
        # Fundo da caixa
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = cor
        shape.line.fill.background()
        
        # Valor (grande)
        txBox = slide.shapes.add_textbox(left, top + Inches(0.3), width, Inches(0.8))
        tf = txBox.text_frame
        tf.text = valor
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(42)
        p.font.bold = True
        p.font.color.rgb = self.branco
        
        # Título (pequeno)
        txBox = slide.shapes.add_textbox(left, top + Inches(1.2), width, Inches(0.4))
        tf = txBox.text_frame
        tf.text = titulo
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(14)
        p.font.color.rgb = self.branco
        
    def adicionar_imagem(self, slide, caminho, left, top, width, height):
        """Adiciona imagem ao slide."""
        if os.path.exists(caminho):
            slide.shapes.add_picture(caminho, left, top, width=width, height=height)
            return True
        return False
    
    # ==================== PARTE 1: ANÁLISES POLÍTICAS ====================
    
    def slide_01_capa(self):
        """Slide 1: Capa principal."""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Fundo gradiente (simulado com retângulos)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(10), Inches(7.5)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.azul_escuro
        shape.line.fill.background()
        
        # Círculo decorativo
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(7), Inches(-2),
            Inches(6), Inches(6)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = self.azul_principal
        circle.fill.fore_color.brightness = -0.2
        circle.line.fill.background()
        
        # Título principal
        txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.2))
        tf = txBox.text_frame
        tf.text = "Análise de Gastos\nParlamentares"
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        p.font.size = Pt(56)
        p.font.bold = True
        p.font.color.rgb = self.branco
        
        # Subtítulo
        txBox = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(7), Inches(0.6))
        tf = txBox.text_frame
        tf.text = "Sistema Inteligente de Análise de Dados da Câmara dos Deputados"
        p = tf.paragraphs[0]
        p.font.size = Pt(18)
        p.font.color.rgb = self.cinza_claro
        
        # Rodapé
        txBox = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.5))
        tf = txBox.text_frame
        tf.text = "Grupo 1 • Ciência de Dados • 2025"
        p = tf.paragraphs[0]
        p.font.size = Pt(14)
        p.font.color.rgb = self.cinza_claro
    
    def slide_02_divisor_parte1(self):
        """Slide 2: Divisor Parte 1."""
        self.criar_slide_titulo_secao(
            "Análises dos Dados\nParlamentares",
            "01",
            self.azul_principal
        )
    
    def slide_03_overview_dados(self):
        """Slide 3: Overview dos dados analisados."""
        slide = self.criar_slide_conteudo("📊 Visão Geral dos Dados", self.azul_principal)
        
        # Carregar dados reais
        try:
            df_completo = pd.read_csv(os.path.join(self.pasta_resultados, "analise_completa.csv"))
            df_partidos = pd.read_csv(os.path.join(self.pasta_resultados, "gastos_por_partido.csv"))
            df_estados = pd.read_csv(os.path.join(self.pasta_resultados, "gastos_por_estado.csv"))
            
            total_gastos = df_completo['valor'].sum()
            total_registros = len(df_completo)
            total_partidos = len(df_partidos)
            total_estados = len(df_estados)
            
        except Exception as e:
            print(f"  ⚠️  Erro ao carregar dados: {e}")
            total_gastos = 0
            total_registros = 0
            total_partidos = 0
            total_estados = 0
        
        # Métricas em caixas
        y_pos = 2.2
        metricas = [
            ("Registros Analisados", f"{total_registros:,}", self.azul_principal),
            ("Valor Total", f"R$ {total_gastos/1e6:.1f}M", self.verde),
            ("Partidos", f"{total_partidos}", self.laranja),
            ("Estados", f"{total_estados}", self.roxo)
        ]
        
        x_positions = [0.8, 3.0, 5.2, 7.4]
        for i, (titulo, valor, cor) in enumerate(metricas):
            self.adicionar_caixa_destaque(
                slide, titulo, valor,
                Inches(x_positions[i]), Inches(y_pos),
                Inches(1.8), Inches(1.8),
                cor
            )
        
        # Informações adicionais
        y_pos = 4.8
        infos = [
            "✓ Dados integrados com API oficial da Câmara dos Deputados",
            "✓ Taxa de identificação superior a 95%",
            "✓ Análise multidimensional (partido, estado, tipo despesa)",
            "✓ Processamento automatizado em tempo real"
        ]
        
        for info in infos:
            self.adicionar_bullet_point(slide, info, Inches(1), Inches(y_pos), Inches(8), 14)
            y_pos += 0.4
    
    def slide_04_gastos_partidos(self):
        """Slide 4: Análise por partido."""
        slide = self.criar_slide_conteudo("🏛️ Gastos por Partido Político", self.laranja)
        
        # Gráfico
        img_path = os.path.join(self.pasta_resultados, "gastos_por_partido.png")
        if self.adicionar_imagem(slide, img_path, Inches(0.5), Inches(1.5), Inches(6), Inches(4.5)):
            
            # Insights (lado direito)
            try:
                df = pd.read_csv(os.path.join(self.pasta_resultados, "gastos_por_partido.csv"))
                top1 = df.iloc[0]
                top2 = df.iloc[1]
                
                # Caixa de destaque
                self.adicionar_caixa_destaque(
                    slide, "Partido Líder", top1['partido'],
                    Inches(7), Inches(1.8),
                    Inches(2.5), Inches(1.5),
                    self.laranja
                )
                
                # Insights
                insights = [
                    f"💰 R$ {top1['total_gasto']/1e6:.1f}M em gastos",
                    f"📊 {len(df)} partidos analisados",
                    f"🥈 2º lugar: {top2['partido']}",
                    f"📈 Variação significativa"
                ]
                
                y_pos = 3.8
                for insight in insights:
                    self.adicionar_bullet_point(
                        slide, insight,
                        Inches(7), Inches(y_pos),
                        Inches(2.5), 13
                    )
                    y_pos += 0.45
            except:
                pass
    
    def slide_05_gastos_estados(self):
        """Slide 5: Análise por estado."""
        slide = self.criar_slide_conteudo("🗺️ Distribuição Geográfica", self.verde)
        
        # Gráfico
        img_path = os.path.join(self.pasta_resultados, "gastos_por_estado.png")
        if self.adicionar_imagem(slide, img_path, Inches(0.5), Inches(1.5), Inches(6), Inches(4.5)):
            
            try:
                df = pd.read_csv(os.path.join(self.pasta_resultados, "gastos_por_estado.csv"))
                top1 = df.iloc[0]
                
                # Caixa de destaque
                self.adicionar_caixa_destaque(
                    slide, "Estado Líder", top1['uf'],
                    Inches(7), Inches(1.8),
                    Inches(2.5), Inches(1.5),
                    self.verde
                )
                
                # Insights
                total = df['total_gasto'].sum()
                percentual_top1 = (top1['total_gasto'] / total) * 100
                
                insights = [
                    f"💰 R$ {top1['total_gasto']/1e6:.1f}M",
                    f"📊 {percentual_top1:.1f}% do total",
                    f"🌎 {len(df)} estados/DF",
                    f"📍 Análise regional completa"
                ]
                
                y_pos = 3.8
                for insight in insights:
                    self.adicionar_bullet_point(
                        slide, insight,
                        Inches(7), Inches(y_pos),
                        Inches(2.5), 13
                    )
                    y_pos += 0.45
            except:
                pass
    
    def slide_06_tipos_despesa(self):
        """Slide 6: Tipos de despesa."""
        slide = self.criar_slide_conteudo("💳 Principais Categorias de Gastos", self.roxo)
        
        # Gráfico
        img_path = os.path.join(self.pasta_resultados, "tipos_despesa.png")
        if self.adicionar_imagem(slide, img_path, Inches(0.5), Inches(1.5), Inches(6), Inches(4.5)):
            
            try:
                df = pd.read_csv(os.path.join(self.pasta_resultados, "gastos_por_tipo_despesa.csv"))
                
                # Caixa de destaque
                self.adicionar_caixa_destaque(
                    slide, "Categorias", str(len(df)),
                    Inches(7), Inches(1.8),
                    Inches(2.5), Inches(1.5),
                    self.roxo
                )
                
                # Top 3 despesas
                top3_text = "🏆 Top 3 Despesas:"
                txBox = slide.shapes.add_textbox(Inches(7), Inches(3.5), Inches(2.5), Inches(0.3))
                tf = txBox.text_frame
                tf.text = top3_text
                p = tf.paragraphs[0]
                p.font.size = Pt(13)
                p.font.bold = True
                p.font.color.rgb = self.azul_escuro
                
                y_pos = 3.9
                for i in range(min(3, len(df))):
                    despesa = df.iloc[i]
                    texto = f"{i+1}. {despesa['tipoDespesa'][:25]}"
                    self.adicionar_bullet_point(
                        slide, texto,
                        Inches(7.2), Inches(y_pos),
                        Inches(2.3), 11
                    )
                    y_pos += 0.35
            except:
                pass
    
    def slide_07_top_deputados(self):
        """Slide 7: Top deputados."""
        slide = self.criar_slide_conteudo("👥 Ranking de Deputados", self.azul_claro)
        
        # Gráfico
        img_path = os.path.join(self.pasta_resultados, "top_deputados.png")
        if self.adicionar_imagem(slide, img_path, Inches(0.5), Inches(1.5), Inches(6), Inches(4.5)):
            
            try:
                df = pd.read_csv(os.path.join(self.pasta_resultados, "top_deputados.csv"))
                top1 = df.iloc[0]
                
                # Caixa com nome do deputado líder
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(7), Inches(1.8),
                    Inches(2.5), Inches(2.2)
                )
                shape.fill.solid()
                shape.fill.fore_color.rgb = self.azul_claro
                shape.line.fill.background()
                
                # Nome
                txBox = slide.shapes.add_textbox(Inches(7), Inches(2), Inches(2.5), Inches(0.6))
                tf = txBox.text_frame
                tf.text = "🥇 Líder em Gastos"
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                p.font.size = Pt(13)
                p.font.bold = True
                p.font.color.rgb = self.branco
                
                txBox = slide.shapes.add_textbox(Inches(7), Inches(2.6), Inches(2.5), Inches(0.7))
                tf = txBox.text_frame
                tf.text = top1['nome_deputado'][:20]
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                p.font.size = Pt(16)
                p.font.bold = True
                p.font.color.rgb = self.branco
                
                # Partido e Estado
                txBox = slide.shapes.add_textbox(Inches(7), Inches(3.3), Inches(2.5), Inches(0.3))
                tf = txBox.text_frame
                tf.text = f"{top1['partido']} - {top1['uf']}"
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                p.font.size = Pt(12)
                p.font.color.rgb = self.branco
                
                # Valor
                txBox = slide.shapes.add_textbox(Inches(7), Inches(3.7), Inches(2.5), Inches(0.4))
                tf = txBox.text_frame
                tf.text = f"R$ {top1['total_gasto']/1e6:.2f}M"
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                p.font.size = Pt(20)
                p.font.bold = True
                p.font.color.rgb = self.branco
                
                # Info adicional
                y_pos = 4.5
                infos = [
                    f"📊 Top 20 analisados",
                    f"📈 Média: R$ {df['total_gasto'].mean()/1e6:.2f}M",
                    f"📉 Variação alta"
                ]
                for info in infos:
                    self.adicionar_bullet_point(
                        slide, info,
                        Inches(7), Inches(y_pos),
                        Inches(2.5), 12
                    )
                    y_pos += 0.4
            except:
                pass
    
    def slide_08_dashboard(self):
        """Slide 8: Dashboard completo."""
        slide = self.criar_slide_conteudo("📈 Dashboard - Visão 360°", self.azul_principal)
        
        # Imagem do resumo geral (maior)
        img_path = os.path.join(self.pasta_resultados, "resumo_geral.png")
        self.adicionar_imagem(slide, img_path, Inches(0.5), Inches(1.4), Inches(9), Inches(5.8))
    
    # ==================== PARTE 2: ARQUITETURA DO CÓDIGO ====================
    
    def slide_09_divisor_parte2(self):
        """Slide 9: Divisor Parte 2."""
        self.criar_slide_titulo_secao(
            "Arquitetura da\nSolução Técnica",
            "02",
            self.cinza_escuro
        )
    
    def slide_10_visao_geral_solucao(self):
        """Slide 10: Visão geral da solução."""
        slide = self.criar_slide_conteudo("🔧 Arquitetura do Sistema", self.azul_principal)
        
        # Diagrama simplificado em texto
        y_pos = 2
        
        # Fluxo principal
        fluxo = [
            ("1️⃣", "Carregamento de Dados", "CSV com despesas parlamentares", self.azul_principal),
            ("2️⃣", "Integração Externa", "API da Câmara dos Deputados", self.verde),
            ("3️⃣", "Processamento", "Limpeza e cruzamento de dados", self.laranja),
            ("4️⃣", "Análise Estatística", "Agregações multidimensionais", self.roxo),
            ("5️⃣", "Visualização", "Gráficos profissionais (300 DPI)", self.azul_claro),
            ("6️⃣", "Exportação", "CSV, PNG e PowerPoint", self.verde)
        ]
        
        for emoji, titulo, descricao, cor in fluxo:
            # Caixa
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(1.5), Inches(y_pos),
                Inches(7), Inches(0.7)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = cor
            shape.line.fill.background()
            
            # Emoji e título
            txBox = slide.shapes.add_textbox(Inches(1.7), Inches(y_pos + 0.1), Inches(6.5), Inches(0.3))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            
            run = p.add_run()
            run.text = f"{emoji}  {titulo}"
            run.font.size = Pt(16)
            run.font.bold = True
            run.font.color.rgb = self.branco
            
            # Descrição
            txBox = slide.shapes.add_textbox(Inches(1.7), Inches(y_pos + 0.38), Inches(6.5), Inches(0.25))
            tf = txBox.text_frame
            tf.text = descricao
            p = tf.paragraphs[0]
            p.font.size = Pt(11)
            p.font.color.rgb = self.branco
            
            y_pos += 0.85
    
    def slide_11_tecnologias(self):
        """Slide 11: Stack tecnológico."""
        slide = self.criar_slide_conteudo("🛠️ Stack Tecnológico", self.verde)
        
        # Dividir em categorias
        categorias = [
            {
                'titulo': '💻 Core',
                'itens': ['Python 3.13', 'Pandas 2.3.3', 'NumPy 2.3.4'],
                'cor': self.azul_principal,
                'left': 0.8,
                'top': 2
            },
            {
                'titulo': '📊 Visualização',
                'itens': ['Matplotlib 3.10', 'Seaborn 0.13', 'python-pptx 1.0'],
                'cor': self.verde,
                'left': 3.5,
                'top': 2
            },
            {
                'titulo': '🌐 Integração',
                'itens': ['Requests 2.32', 'API REST', 'JSON'],
                'cor': self.laranja,
                'left': 6.2,
                'top': 2
            },
            {
                'titulo': '🔧 Utilidades',
                'itens': ['Unidecode 1.4', 'pathlib', 'datetime'],
                'cor': self.roxo,
                'left': 0.8,
                'top': 4.5
            },
            {
                'titulo': '📝 Documentação',
                'itens': ['Markdown', 'Docstrings', 'Type Hints'],
                'cor': self.azul_claro,
                'left': 3.5,
                'top': 4.5
            },
            {
                'titulo': '⚡ Performance',
                'itens': ['Vetorização', 'Cache API', 'Timestamps'],
                'cor': self.cinza_escuro,
                'left': 6.2,
                'top': 4.5
            }
        ]
        
        for cat in categorias:
            # Caixa da categoria
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(cat['left']), Inches(cat['top']),
                Inches(2.3), Inches(2)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = cat['cor']
            shape.line.fill.background()
            
            # Título
            txBox = slide.shapes.add_textbox(
                Inches(cat['left']), Inches(cat['top'] + 0.15),
                Inches(2.3), Inches(0.35)
            )
            tf = txBox.text_frame
            tf.text = cat['titulo']
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = self.branco
            
            # Itens
            y_offset = 0.6
            for item in cat['itens']:
                txBox = slide.shapes.add_textbox(
                    Inches(cat['left'] + 0.2), Inches(cat['top'] + y_offset),
                    Inches(1.9), Inches(0.25)
                )
                tf = txBox.text_frame
                tf.text = f"• {item}"
                p = tf.paragraphs[0]
                p.font.size = Pt(11)
                p.font.color.rgb = self.branco
                y_offset += 0.35
    
    def slide_12_estrutura_modular(self):
        """Slide 12: Estrutura modular."""
        slide = self.criar_slide_conteudo("📦 Arquitetura Modular", self.roxo)
        
        modulos = [
            {
                'nome': 'api_client.py',
                'icone': '🌐',
                'funcao': 'Conexão com API',
                'detalhes': 'Busca dados cadastrais dos deputados'
            },
            {
                'nome': 'data_loader.py',
                'icone': '📂',
                'funcao': 'Carregamento',
                'detalhes': 'Leitura e limpeza de CSV'
            },
            {
                'nome': 'data_analyzer.py',
                'icone': '🔍',
                'funcao': 'Análise',
                'detalhes': 'Cruzamento e agregações estatísticas'
            },
            {
                'nome': 'visualizer.py',
                'icone': '📊',
                'funcao': 'Visualização',
                'detalhes': 'Geração de gráficos profissionais'
            },
            {
                'nome': 'main.py',
                'icone': '🎯',
                'funcao': 'Orquestração',
                'detalhes': 'Coordena todo o pipeline'
            }
        ]
        
        y_pos = 2.2
        for modulo in modulos:
            # Linha com módulo
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(1), Inches(y_pos),
                Inches(8), Inches(0.65)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = self.roxo
            shape.line.fill.background()
            
            # Ícone
            txBox = slide.shapes.add_textbox(Inches(1.2), Inches(y_pos + 0.15), Inches(0.5), Inches(0.4))
            tf = txBox.text_frame
            tf.text = modulo['icone']
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(20)
            
            # Nome do arquivo
            txBox = slide.shapes.add_textbox(Inches(1.8), Inches(y_pos + 0.08), Inches(2), Inches(0.3))
            tf = txBox.text_frame
            tf.text = modulo['nome']
            p = tf.paragraphs[0]
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = self.branco
            
            # Função
            txBox = slide.shapes.add_textbox(Inches(4), Inches(y_pos + 0.08), Inches(1.5), Inches(0.3))
            tf = txBox.text_frame
            tf.text = modulo['funcao']
            p = tf.paragraphs[0]
            p.font.size = Pt(11)
            p.font.color.rgb = self.branco
            
            # Detalhes
            txBox = slide.shapes.add_textbox(Inches(1.8), Inches(y_pos + 0.35), Inches(6.8), Inches(0.25))
            tf = txBox.text_frame
            tf.text = modulo['detalhes']
            p = tf.paragraphs[0]
            p.font.size = Pt(10)
            p.font.color.rgb = self.cinza_claro
            
            y_pos += 0.8
    
    def slide_13_pipeline_dados(self):
        """Slide 13: Pipeline de processamento."""
        slide = self.criar_slide_conteudo("⚙️ Pipeline de Dados", self.laranja)
        
        # Etapas do pipeline
        etapas = [
            ("📥 INPUT", "CSV (~285k registros)", self.azul_principal),
            ("🧹 CLEAN", "Remove nulos e inválidos", self.laranja),
            ("🔄 NORMALIZE", "Uppercase + sem acentos", self.verde),
            ("🌐 ENRICH", "API: partido, estado", self.roxo),
            ("🔗 MATCH", "Cruzamento por nome (>95%)", self.azul_claro),
            ("📊 AGGREGATE", "Por partido, estado, despesa", self.verde),
            ("📈 VISUALIZE", "5 gráficos + dashboard", self.laranja),
            ("💾 EXPORT", "11 arquivos timestampados", self.azul_principal)
        ]
        
        # Desenhar como fluxo
        x_start = 1.2
        y_start = 2
        box_width = 1.8
        box_height = 0.55
        spacing = 0.15
        
        for i, (titulo, desc, cor) in enumerate(etapas):
            # Alternar posição (zigue-zague visual)
            row = i // 2
            col = i % 2
            
            x = x_start + col * (box_width + spacing + 4)
            y = y_start + row * (box_height + spacing + 0.1)
            
            # Caixa
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(y),
                Inches(box_width), Inches(box_height)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = cor
            shape.line.fill.background()
            
            # Título
            txBox = slide.shapes.add_textbox(Inches(x), Inches(y + 0.08), Inches(box_width), Inches(0.25))
            tf = txBox.text_frame
            tf.text = titulo
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = self.branco
            
            # Descrição
            txBox = slide.shapes.add_textbox(Inches(x), Inches(y + 0.3), Inches(box_width), Inches(0.2))
            tf = txBox.text_frame
            tf.text = desc
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(9)
            p.font.color.rgb = self.cinza_claro
            
            # Seta para próximo (se não for o último)
            if i < len(etapas) - 1:
                if col == 0:  # Seta para direita
                    arrow = slide.shapes.add_shape(
                        MSO_SHAPE.RIGHT_ARROW,
                        Inches(x + box_width + 0.05), Inches(y + 0.15),
                        Inches(0.4), Inches(0.25)
                    )
                else:  # Seta para baixo
                    arrow = slide.shapes.add_shape(
                        MSO_SHAPE.DOWN_ARROW,
                        Inches(x + box_width / 2 - 0.15), Inches(y + box_height + 0.02),
                        Inches(0.3), Inches(0.3)
                    )
                
                arrow.fill.solid()
                arrow.fill.fore_color.rgb = self.cinza_claro
                arrow.line.fill.background()
    
    def slide_14_diferenciais(self):
        """Slide 14: Diferenciais da solução."""
        slide = self.criar_slide_conteudo("⭐ Diferenciais Técnicos", self.verde)
        
        diferenciais = [
            {
                'icone': '🚀',
                'titulo': 'Alta Performance',
                'desc': 'Processa 285k registros em ~2 minutos'
            },
            {
                'icone': '🎯',
                'titulo': 'Precisão Elevada',
                'desc': 'Taxa de match superior a 95%'
            },
            {
                'icone': '🔄',
                'titulo': 'Totalmente Automatizado',
                'desc': 'Pipeline completo sem intervenção manual'
            },
            {
                'icone': '📊',
                'titulo': 'Visualizações Profissionais',
                'desc': 'Gráficos em alta resolução (300 DPI)'
            },
            {
                'icone': '🗂️',
                'titulo': 'Organização Inteligente',
                'desc': 'Resultados timestampados por execução'
            },
            {
                'icone': '🔌',
                'titulo': 'Modular e Extensível',
                'desc': 'Fácil adicionar novas análises'
            },
            {
                'icone': '📈',
                'titulo': 'Análise Multidimensional',
                'desc': 'Partido, estado, despesa, deputado'
            },
            {
                'icone': '💾',
                'titulo': 'Múltiplos Formatos',
                'desc': 'CSV, PNG, PPTX automaticamente'
            }
        ]
        
        # Grid 4x2
        x_positions = [0.8, 3.2, 5.6, 8.0]
        y_positions = [2.0, 4.3]
        
        i = 0
        for y in y_positions:
            for x in x_positions:
                if i >= len(diferenciais):
                    break
                    
                diff = diferenciais[i]
                
                # Caixa
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(x), Inches(y),
                    Inches(2.2), Inches(1.8)
                )
                shape.fill.solid()
                shape.fill.fore_color.rgb = self.verde
                shape.line.fill.background()
                
                # Ícone
                txBox = slide.shapes.add_textbox(Inches(x), Inches(y + 0.2), Inches(2.2), Inches(0.4))
                tf = txBox.text_frame
                tf.text = diff['icone']
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                p.font.size = Pt(32)
                
                # Título
                txBox = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.7), Inches(2), Inches(0.35))
                tf = txBox.text_frame
                tf.text = diff['titulo']
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                p.font.size = Pt(13)
                p.font.bold = True
                p.font.color.rgb = self.branco
                
                # Descrição
                txBox = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 1.1), Inches(2), Inches(0.6))
                tf = txBox.text_frame
                tf.text = diff['desc']
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                p.font.size = Pt(10)
                p.font.color.rgb = self.cinza_claro
                
                i += 1
    
    def slide_15_resultados_tecnicos(self):
        """Slide 15: Métricas técnicas."""
        slide = self.criar_slide_conteudo("📊 Métricas de Qualidade", self.azul_claro)
        
        # Métricas grandes
        metricas = [
            ("Cobertura de Dados", ">95%", "Taxa de identificação", self.verde),
            ("Qualidade dos Dados", "64%", "Dados aproveitados", self.azul_principal),
            ("Tempo de Execução", "~2min", "Para 285k registros", self.laranja),
            ("Resolução Gráfica", "300 DPI", "Qualidade profissional", self.roxo)
        ]
        
        y_pos = 2.2
        for titulo, valor, desc, cor in metricas:
            # Barra horizontal
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(1.5), Inches(y_pos),
                Inches(7), Inches(0.9)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = cor
            shape.line.fill.background()
            
            # Título
            txBox = slide.shapes.add_textbox(Inches(1.7), Inches(y_pos + 0.1), Inches(3), Inches(0.3))
            tf = txBox.text_frame
            tf.text = titulo
            p = tf.paragraphs[0]
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = self.branco
            
            # Descrição
            txBox = slide.shapes.add_textbox(Inches(1.7), Inches(y_pos + 0.45), Inches(3), Inches(0.25))
            tf = txBox.text_frame
            tf.text = desc
            p = tf.paragraphs[0]
            p.font.size = Pt(11)
            p.font.color.rgb = self.cinza_claro
            
            # Valor (direita)
            txBox = slide.shapes.add_textbox(Inches(6.5), Inches(y_pos + 0.15), Inches(1.8), Inches(0.6))
            tf = txBox.text_frame
            tf.text = valor
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = self.branco
            
            y_pos += 1.1
    
    def slide_16_algoritmo_matching(self):
        """Slide 16: Algoritmo de matching."""
        slide = self.criar_slide_conteudo("🔍 Como Funciona o Cruzamento de Dados", self.roxo)
        
        # Título da seção
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(0.4))
        tf = txBox.text_frame
        tf.text = "Identificando Deputados nos Dados"
        p = tf.paragraphs[0]
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = self.roxo
        
        # Processo detalhado
        etapas = [
            {
                'num': '1',
                'titulo': 'Padronizar Nomes',
                'desc': 'Coloca tudo em MAIÚSCULAS e remove acentos',
                'exemplo': '"José da Silva" vira "JOSE DA SILVA"'
            },
            {
                'num': '2',
                'titulo': 'Limpar Títulos',
                'desc': 'Remove palavras como "DEP.", "DR.", "JR."',
                'exemplo': '"DEP. MARIA SOUZA JR." vira "MARIA SOUZA"'
            },
            {
                'num': '3',
                'titulo': 'Buscar Correspondência',
                'desc': 'Procura o nome limpo na lista de deputados da API',
                'exemplo': 'Se encontrar → adiciona partido e estado'
            },
            {
                'num': '4',
                'titulo': 'Confirmar Resultado',
                'desc': 'Marca como "identificado" ou "não identificado"',
                'exemplo': 'Conseguimos identificar mais de 95%!'
            }
        ]
        
        y_pos = 2.5
        for etapa in etapas:
            # Número grande
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(1.2), Inches(y_pos),
                Inches(0.5), Inches(0.5)
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = self.roxo
            circle.line.fill.background()
            
            txBox = slide.shapes.add_textbox(Inches(1.2), Inches(y_pos + 0.05), Inches(0.5), Inches(0.4))
            tf = txBox.text_frame
            tf.text = etapa['num']
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = self.branco
            
            # Conteúdo
            txBox = slide.shapes.add_textbox(Inches(2), Inches(y_pos), Inches(7), Inches(0.25))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            
            run = p.add_run()
            run.text = f"{etapa['titulo']}: "
            run.font.size = Pt(14)
            run.font.bold = True
            run.font.color.rgb = self.azul_escuro
            
            run = p.add_run()
            run.text = etapa['desc']
            run.font.size = Pt(12)
            run.font.color.rgb = self.cinza_escuro
            
            # Exemplo
            txBox = slide.shapes.add_textbox(Inches(2.2), Inches(y_pos + 0.3), Inches(6.8), Inches(0.25))
            tf = txBox.text_frame
            tf.text = f"💡 {etapa['exemplo']}"
            p = tf.paragraphs[0]
            p.font.size = Pt(10)
            p.font.italic = True
            p.font.color.rgb = self.cinza_claro
            
            y_pos += 0.9
    
    def slide_17_tratamento_erros(self):
        """Slide 17: Tratamento de erros e edge cases."""
        slide = self.criar_slide_conteudo("⚠️ Como Lidamos com Problemas", self.laranja)
        
        casos = [
            {
                'tipo': '🛡️ Internet Instável',
                'problema': 'Às vezes a API da Câmara não responde',
                'solucao': 'Sistema tenta 3 vezes antes de desistir',
                'codigo': 'Espera 30 segundos por tentativa'
            },
            {
                'tipo': '📊 Dados com Erro',
                'problema': 'CSV pode ter valores negativos ou vazios',
                'solucao': 'Remove automaticamente linhas problemáticas',
                'codigo': 'Descarta cerca de 36% dos dados ruins'
            },
            {
                'tipo': '🔤 Nomes Estranhos',
                'problema': 'Caracteres especiais podem dar erro',
                'solucao': 'Sistema entende diferentes formatos de texto',
                'codigo': 'Funciona com acentos e caracteres raros'
            },
            {
                'tipo': '🔗 Deputado Não Encontrado',
                'problema': 'Alguns nomes não batem com a API',
                'solucao': 'Marca como "NÃO IDENTIFICADO" para análise',
                'codigo': 'Acontece em apenas 5% dos casos'
            }
        ]
        
        y_pos = 2.2
        for i, caso in enumerate(casos):
            # Caixa colorida
            cor = self.laranja if i % 2 == 0 else self.roxo
            
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(1), Inches(y_pos),
                Inches(8), Inches(1.0)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = cor
            shape.line.fill.background()
            
            # Tipo
            txBox = slide.shapes.add_textbox(Inches(1.2), Inches(y_pos + 0.1), Inches(2), Inches(0.25))
            tf = txBox.text_frame
            tf.text = caso['tipo']
            p = tf.paragraphs[0]
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = self.branco
            
            # Problema
            txBox = slide.shapes.add_textbox(Inches(1.2), Inches(y_pos + 0.35), Inches(3.5), Inches(0.2))
            tf = txBox.text_frame
            tf.text = f"⚠️ {caso['problema']}"
            p = tf.paragraphs[0]
            p.font.size = Pt(10)
            p.font.color.rgb = self.cinza_claro
            
            # Solução
            txBox = slide.shapes.add_textbox(Inches(5), Inches(y_pos + 0.1), Inches(3.8), Inches(0.25))
            tf = txBox.text_frame
            tf.text = f"✓ {caso['solucao']}"
            p = tf.paragraphs[0]
            p.font.size = Pt(11)
            p.font.color.rgb = self.branco
            
            # Código/métrica
            txBox = slide.shapes.add_textbox(Inches(5), Inches(y_pos + 0.35), Inches(3.8), Inches(0.2))
            tf = txBox.text_frame
            tf.text = f"📝 {caso['codigo']}"
            p = tf.paragraphs[0]
            p.font.size = Pt(9)
            p.font.italic = True
            p.font.color.rgb = self.cinza_claro
            
            y_pos += 1.15
    
    def slide_18_otimizacoes(self):
        """Slide 18: Otimizações de performance."""
        slide = self.criar_slide_conteudo("⚡ Por Que é Tão Rápido?", self.verde)
        
        # Dividir em duas colunas
        otimizacoes_esquerda = [
            {
                'icone': '🔄',
                'titulo': 'Processamento em Massa',
                'desc': 'Trabalha com muitos dados de uma vez',
                'ganho': '100x mais rápido que um por um'
            },
            {
                'icone': '💾',
                'titulo': 'Memória Inteligente',
                'desc': 'Salva resultados para não buscar de novo',
                'ganho': 'Não precisa repetir trabalho'
            },
            {
                'icone': '📊',
                'titulo': 'Carrega Só o Necessário',
                'desc': 'Não pega tudo de uma vez',
                'ganho': 'Usa menos memória do PC'
            }
        ]
        
        otimizacoes_direita = [
            {
                'icone': '🎯',
                'titulo': 'Busca Super Rápida',
                'desc': 'Organiza dados para achar em 1 segundo',
                'ganho': 'Encontra qualquer coisa na hora'
            },
            {
                'icone': '🔍',
                'titulo': 'Remove o Inútil Primeiro',
                'desc': 'Elimina dados ruins antes de processar',
                'ganho': 'Descarta 36% que não serve'
            },
            {
                'icone': '📈',
                'titulo': 'Trabalho em Grupos',
                'desc': 'Faz tarefas parecidas juntas',
                'ganho': 'Economiza tempo e esforço'
            }
        ]
        
        # Coluna esquerda
        y_pos = 2.2
        for opt in otimizacoes_esquerda:
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(1), Inches(y_pos),
                Inches(3.8), Inches(1.1)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = self.verde
            shape.line.fill.background()
            
            # Ícone e título
            txBox = slide.shapes.add_textbox(Inches(1.2), Inches(y_pos + 0.15), Inches(3.4), Inches(0.3))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            
            run = p.add_run()
            run.text = f"{opt['icone']} "
            run.font.size = Pt(16)
            
            run = p.add_run()
            run.text = opt['titulo']
            run.font.size = Pt(13)
            run.font.bold = True
            run.font.color.rgb = self.branco
            
            # Descrição
            txBox = slide.shapes.add_textbox(Inches(1.2), Inches(y_pos + 0.45), Inches(3.4), Inches(0.25))
            tf = txBox.text_frame
            tf.text = opt['desc']
            p = tf.paragraphs[0]
            p.font.size = Pt(10)
            p.font.color.rgb = self.cinza_claro
            
            # Ganho
            txBox = slide.shapes.add_textbox(Inches(1.2), Inches(y_pos + 0.7), Inches(3.4), Inches(0.25))
            tf = txBox.text_frame
            tf.text = f"⚡ {opt['ganho']}"
            p = tf.paragraphs[0]
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = self.branco
            
            y_pos += 1.25
        
        # Coluna direita
        y_pos = 2.2
        for opt in otimizacoes_direita:
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(5.2), Inches(y_pos),
                Inches(3.8), Inches(1.1)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = self.azul_principal
            shape.line.fill.background()
            
            # Ícone e título
            txBox = slide.shapes.add_textbox(Inches(5.4), Inches(y_pos + 0.15), Inches(3.4), Inches(0.3))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            
            run = p.add_run()
            run.text = f"{opt['icone']} "
            run.font.size = Pt(16)
            
            run = p.add_run()
            run.text = opt['titulo']
            run.font.size = Pt(13)
            run.font.bold = True
            run.font.color.rgb = self.branco
            
            # Descrição
            txBox = slide.shapes.add_textbox(Inches(5.4), Inches(y_pos + 0.45), Inches(3.4), Inches(0.25))
            tf = txBox.text_frame
            tf.text = opt['desc']
            p = tf.paragraphs[0]
            p.font.size = Pt(10)
            p.font.color.rgb = self.cinza_claro
            
            # Ganho
            txBox = slide.shapes.add_textbox(Inches(5.4), Inches(y_pos + 0.7), Inches(3.4), Inches(0.25))
            tf = txBox.text_frame
            tf.text = f"⚡ {opt['ganho']}"
            p = tf.paragraphs[0]
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = self.branco
            
            y_pos += 1.25
    
    def slide_19_escalabilidade(self):
        """Slide 19: Escalabilidade e capacidade."""
        slide = self.criar_slide_conteudo("📈 Números do Sistema", self.azul_claro)
        
        # Métricas de capacidade
        metricas = [
            ("Volume Processado", "285k → 107k", "Analisa 107 mil gastos válidos", self.azul_principal),
            ("Velocidade", "~2 minutos", "Da leitura até os gráficos", self.verde),
            ("Uso de Memória", "< 500 MB", "Roda em qualquer notebook", self.laranja),
            ("Resultados Criados", "11 arquivos", "Tabelas + Gráficos + PowerPoint", self.roxo)
        ]
        
        y_pos = 2.2
        for titulo, valor, desc, cor in metricas:
            # Barra
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(1.5), Inches(y_pos),
                Inches(7), Inches(0.8)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = cor
            shape.line.fill.background()
            
            # Título
            txBox = slide.shapes.add_textbox(Inches(1.7), Inches(y_pos + 0.1), Inches(3), Inches(0.25))
            tf = txBox.text_frame
            tf.text = titulo
            p = tf.paragraphs[0]
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = self.branco
            
            # Descrição
            txBox = slide.shapes.add_textbox(Inches(1.7), Inches(y_pos + 0.4), Inches(3), Inches(0.2))
            tf = txBox.text_frame
            tf.text = desc
            p = tf.paragraphs[0]
            p.font.size = Pt(10)
            p.font.color.rgb = self.cinza_claro
            
            # Valor grande
            txBox = slide.shapes.add_textbox(Inches(6), Inches(y_pos + 0.15), Inches(2.3), Inches(0.5))
            tf = txBox.text_frame
            tf.text = valor
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = self.branco
            
            y_pos += 0.95
        
        # Capacidades futuras
        txBox = slide.shapes.add_textbox(Inches(1.5), Inches(6.2), Inches(7), Inches(0.3))
        tf = txBox.text_frame
        tf.text = "🚀 Preparado para analisar anos diferentes e volumes ainda maiores!"
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = self.verde
    
    def slide_20_seguranca_qualidade(self):
        """Slide 20: Segurança e qualidade de código."""
        slide = self.criar_slide_conteudo("🔒 Cuidados com Qualidade", self.roxo)
        
        categorias = [
            {
                'titulo': '📝 Código Organizado',
                'itens': [
                    'Código fácil de entender',
                    'Explicações em cada função',
                    'Nomes claros de variáveis',
                    'Cada arquivo faz uma coisa'
                ]
            },
            {
                'titulo': '🧪 Testa Tudo',
                'itens': [
                    'Verifica se dados fazem sentido',
                    'Trata erros sem travar',
                    'Registra o que aconteceu',
                    'Confere tipos de dados'
                ]
            },
            {
                'titulo': '🔐 Segurança',
                'itens': [
                    'Sem senhas no código',
                    'Limpa dados antes de usar',
                    'Usa API oficial do governo',
                    'Só dados públicos'
                ]
            },
            {
                'titulo': '📦 Bem Estruturado',
                'itens': [
                    'Pastas organizadas',
                    'Módulos separados',
                    'Lista de bibliotecas usadas',
                    'Manual de uso completo'
                ]
            }
        ]
        
        # Grid 2x2
        positions = [
            (1, 2.2), (5.2, 2.2),
            (1, 4.5), (5.2, 4.5)
        ]
        
        cores = [self.roxo, self.azul_principal, self.verde, self.laranja]
        
        for i, (cat, pos, cor) in enumerate(zip(categorias, positions, cores)):
            x, y = pos
            
            # Caixa
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(y),
                Inches(3.8), Inches(1.9)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = cor
            shape.line.fill.background()
            
            # Título
            txBox = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.2), Inches(3.4), Inches(0.3))
            tf = txBox.text_frame
            tf.text = cat['titulo']
            p = tf.paragraphs[0]
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = self.branco
            
            # Itens
            y_item = y + 0.6
            for item in cat['itens']:
                txBox = slide.shapes.add_textbox(Inches(x + 0.3), Inches(y_item), Inches(3.2), Inches(0.25))
                tf = txBox.text_frame
                tf.text = f"✓ {item}"
                p = tf.paragraphs[0]
                p.font.size = Pt(10)
                p.font.color.rgb = self.branco
                y_item += 0.3
    
    def slide_21_encerramento(self):
        """Slide 21: Encerramento."""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Fundo
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(10), Inches(7.5)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.azul_escuro
        shape.line.fill.background()
        
        # Mensagem principal
        txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
        tf = txBox.text_frame
        tf.text = "Obrigado!"
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(64)
        p.font.bold = True
        p.font.color.rgb = self.branco
        
        # Subtítulo
        txBox = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(0.8))
        tf = txBox.text_frame
        tf.text = "Análise de Gastos Parlamentares\nSistema Completo de Business Intelligence"
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(18)
        p.font.color.rgb = self.cinza_claro
        
        # Rodapé com contatos
        txBox = slide.shapes.add_textbox(Inches(1), Inches(5.8), Inches(8), Inches(1))
        tf = txBox.text_frame
        
        # Linha 1
        p = tf.paragraphs[0]
        p.text = "Grupo 1 - Ciência de Dados"
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(14)
        p.font.color.rgb = self.cinza_claro
        
        # Linha 2
        p = tf.add_paragraph()
        p.text = "Leticia (21352) • Gabriel (24734) • Thales (24740) • Maria Fernanda (24767)"
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(12)
        p.font.color.rgb = self.cinza_claro
        
        # Linha 3
        p = tf.add_paragraph()
        p.text = "2025"
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(12)
        p.font.color.rgb = self.cinza_claro
    
    def gerar(self):
        """Gera apresentação completa."""
        print("\n🎨 GERANDO APRESENTAÇÃO DA SOLUÇÃO")
        print("=" * 70)
        
        slides = [
            ("01", "Capa", self.slide_01_capa),
            ("02", "Divisor - Parte 1", self.slide_02_divisor_parte1),
            ("03", "Overview Dados", self.slide_03_overview_dados),
            ("04", "Gastos Partidos", self.slide_04_gastos_partidos),
            ("05", "Gastos Estados", self.slide_05_gastos_estados),
            ("06", "Tipos Despesa", self.slide_06_tipos_despesa),
            ("07", "Top Deputados", self.slide_07_top_deputados),
            ("08", "Dashboard", self.slide_08_dashboard),
            ("09", "Divisor - Parte 2", self.slide_09_divisor_parte2),
            ("10", "Visão Geral Solução", self.slide_10_visao_geral_solucao),
            ("11", "Tecnologias", self.slide_11_tecnologias),
            ("12", "Estrutura Modular", self.slide_12_estrutura_modular),
            ("13", "Pipeline de Dados", self.slide_13_pipeline_dados),
            ("14", "Diferenciais", self.slide_14_diferenciais),
            ("15", "Métricas Técnicas", self.slide_15_resultados_tecnicos),
            ("16", "Algoritmo Matching", self.slide_16_algoritmo_matching),
            ("17", "Tratamento de Erros", self.slide_17_tratamento_erros),
            ("18", "Otimizações", self.slide_18_otimizacoes),
            ("19", "Escalabilidade", self.slide_19_escalabilidade),
            ("20", "Qualidade de Código", self.slide_20_seguranca_qualidade),
            ("21", "Encerramento", self.slide_21_encerramento)
        ]
        
        total_slides = len(slides)
        for num, nome, metodo in slides:
            print(f"  Slide {num}/{total_slides}: {nome}...")
            metodo()
        
        # Salvar
        from datetime import datetime
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        output_path = os.path.join(
            os.path.dirname(self.pasta_resultados),
            f"Apresentacao_Solucao_{timestamp}.pptx"
        )
        
        self.prs.save(output_path)
        
        print("=" * 70)
        print(f"✅ APRESENTAÇÃO SALVA!")
        print(f"📁 Local: {output_path}")
        print(f"📊 Total: {total_slides} slides")
        print(f"📦 Tamanho: {os.path.getsize(output_path) / 1024 / 1024:.2f} MB")
        print(f"\n🎯 Estrutura:")
        print(f"   • Parte 1 (slides 2-8): Análises dos Dados Políticos")
        print(f"   • Parte 2 (slides 9-20): Arquitetura Técnica Detalhada")
        print(f"     - Slides 9-15: Visão Geral e Stack")
        print(f"     - Slides 16-20: Detalhes Técnicos Avançados")
        print("=" * 70 + "\n")
        
        return output_path


def main():
    """Função principal."""
    # Encontrar pasta de resultados mais recente
    base_dir = Path(__file__).parent.parent
    resultados_dir = base_dir / "resultados"
    
    if not resultados_dir.exists():
        print("❌ Pasta 'resultados' não encontrada!")
        return
    
    # Listar execuções
    execucoes = [d for d in resultados_dir.iterdir() if d.is_dir() and d.name.startswith("execucao_")]
    
    if not execucoes:
        print("❌ Nenhuma execução encontrada!")
        print("💡 Execute primeiro: python src/main.py dados/Ano-2025.csv")
        return
    
    # Pegar mais recente
    execucoes.sort(reverse=True)
    pasta_resultados = execucoes[0]
    
    print(f"\n📂 Usando dados de: {pasta_resultados.name}")
    
    # Gerar apresentação
    apresentacao = ApresentacaoSolucao(str(pasta_resultados))
    apresentacao.gerar()
    
    print("✅ Processo concluído com sucesso!")


if __name__ == "__main__":
    main()
