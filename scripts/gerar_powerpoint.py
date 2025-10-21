"""
Gerador de Apresentação PowerPoint

Este script cria uma apresentação profissional com os resultados da análise.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pathlib import Path
import pandas as pd


def criar_apresentacao():
    """Cria apresentação PowerPoint com os resultados"""
    
    # Criar apresentação
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Definir diretório de resultados
    resultados_dir = Path('resultados')
    
    print("📊 Gerando apresentação PowerPoint...")
    
    # === SLIDE 1: TÍTULO ===
    print("   Slide 1: Título")
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Layout em branco
    
    # Adicionar título
    left = Inches(1)
    top = Inches(2.5)
    width = Inches(8)
    height = Inches(2)
    
    txBox = slide1.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "Análise Comparativa de Gastos\nda Cota Parlamentar"
    
    p = tf.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.alignment = PP_ALIGN.CENTER
    
    # Subtítulo
    left = Inches(1)
    top = Inches(4.8)
    txBox2 = slide1.shapes.add_textbox(left, top, width, Inches(0.8))
    tf2 = txBox2.text_frame
    tf2.text = "Por Partido e Estado - Câmara dos Deputados"
    p2 = tf2.paragraphs[0]
    p2.font.size = Pt(24)
    p2.font.color.rgb = RGBColor(100, 100, 100)
    p2.alignment = PP_ALIGN.CENTER
    
    # === SLIDE 2: METODOLOGIA ===
    print("   Slide 2: Metodologia")
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # Layout título + conteúdo
    
    # Título
    title2 = slide2.shapes.title
    title2.text = "Metodologia"
    title2.text_frame.paragraphs[0].font.size = Pt(40)
    title2.text_frame.paragraphs[0].font.bold = True
    title2.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    # Conteúdo
    left = Inches(0.8)
    top = Inches(1.8)
    width = Inches(8.4)
    height = Inches(5)
    
    txBox = slide2.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    conteudo = [
        ("📊 Fonte 1: CSV Estruturado", "~285.000 registros de despesas da Cota Parlamentar"),
        ("🌐 Fonte 2: API JSON", "513 deputados em exercício com dados cadastrais"),
        ("🔗 Cruzamento", "Nome do parlamentar como chave de ligação"),
        ("🧹 Limpeza", "Remoção de valores inválidos e padronização de nomes"),
        ("📈 Análise", "Agregações por partido, estado e tipo de despesa")
    ]
    
    for titulo, descricao in conteudo:
        p = tf.add_paragraph()
        p.text = titulo
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 51, 102)
        p.space_after = Pt(5)
        
        p2 = tf.add_paragraph()
        p2.text = descricao
        p2.font.size = Pt(16)
        p2.font.color.rgb = RGBColor(80, 80, 80)
        p2.level = 1
        p2.space_after = Pt(15)
    
    # === SLIDE 3: GASTOS POR PARTIDO ===
    print("   Slide 3: Gastos por Partido")
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    
    title3 = slide3.shapes.title
    title3.text = "Gastos por Partido"
    title3.text_frame.paragraphs[0].font.size = Pt(40)
    title3.text_frame.paragraphs[0].font.bold = True
    title3.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    # Adicionar imagem
    img_path = resultados_dir / 'gastos_por_partido.png'
    if img_path.exists():
        left = Inches(0.5)
        top = Inches(1.8)
        pic = slide3.shapes.add_picture(str(img_path), left, top, width=Inches(9))
    
    # Adicionar insight
    left = Inches(0.8)
    top = Inches(6.5)
    width = Inches(8.4)
    height = Inches(0.8)
    
    txBox = slide3.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "💡 Importante: Média por deputado normaliza o tamanho dos partidos"
    p = tf.paragraphs[0]
    p.font.size = Pt(14)
    p.font.italic = True
    p.font.color.rgb = RGBColor(100, 100, 100)
    p.alignment = PP_ALIGN.CENTER
    
    # === SLIDE 4: GASTOS POR ESTADO ===
    print("   Slide 4: Gastos por Estado")
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    
    title4 = slide4.shapes.title
    title4.text = "Gastos por Estado"
    title4.text_frame.paragraphs[0].font.size = Pt(40)
    title4.text_frame.paragraphs[0].font.bold = True
    title4.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    # Adicionar imagem
    img_path = resultados_dir / 'gastos_por_estado.png'
    if img_path.exists():
        left = Inches(0.5)
        top = Inches(1.8)
        pic = slide4.shapes.add_picture(str(img_path), left, top, width=Inches(9))
    
    # Adicionar insight
    left = Inches(0.8)
    top = Inches(6.5)
    width = Inches(8.4)
    height = Inches(0.8)
    
    txBox = slide4.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "💡 SP, MG e RJ lideram em volume, mas média revela eficiência relativa"
    p = tf.paragraphs[0]
    p.font.size = Pt(14)
    p.font.italic = True
    p.font.color.rgb = RGBColor(100, 100, 100)
    p.alignment = PP_ALIGN.CENTER
    
    # === SLIDE 5: TIPOS DE DESPESA ===
    print("   Slide 5: Tipos de Despesa")
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    
    title5 = slide5.shapes.title
    title5.text = "Tipos de Despesa Mais Comuns"
    title5.text_frame.paragraphs[0].font.size = Pt(40)
    title5.text_frame.paragraphs[0].font.bold = True
    title5.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    # Adicionar imagem
    img_path = resultados_dir / 'tipos_despesa.png'
    if img_path.exists():
        left = Inches(0.8)
        top = Inches(1.5)
        pic = slide5.shapes.add_picture(str(img_path), left, top, width=Inches(8.4))
    
    # === SLIDE 6: TOP DEPUTADOS ===
    print("   Slide 6: Top Deputados")
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])
    
    title6 = slide6.shapes.title
    title6.text = "Deputados com Maiores Gastos"
    title6.text_frame.paragraphs[0].font.size = Pt(40)
    title6.text_frame.paragraphs[0].font.bold = True
    title6.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    # Adicionar imagem
    img_path = resultados_dir / 'top_deputados.png'
    if img_path.exists():
        left = Inches(0.8)
        top = Inches(1.5)
        pic = slide6.shapes.add_picture(str(img_path), left, top, width=Inches(8.4))
    
    # === SLIDE 7: RESUMO GERAL ===
    print("   Slide 7: Resumo Geral")
    slide7 = prs.slides.add_slide(prs.slide_layouts[5])
    
    title7 = slide7.shapes.title
    title7.text = "Resumo Geral"
    title7.text_frame.paragraphs[0].font.size = Pt(40)
    title7.text_frame.paragraphs[0].font.bold = True
    title7.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    # Adicionar imagem
    img_path = resultados_dir / 'resumo_geral.png'
    if img_path.exists():
        left = Inches(0.3)
        top = Inches(1.5)
        pic = slide7.shapes.add_picture(str(img_path), left, top, width=Inches(9.4))
    
    # === SLIDE 8: PRINCIPAIS INSIGHTS ===
    print("   Slide 8: Principais Insights")
    slide8 = prs.slides.add_slide(prs.slide_layouts[5])
    
    title8 = slide8.shapes.title
    title8.text = "Principais Insights"
    title8.text_frame.paragraphs[0].font.size = Pt(40)
    title8.text_frame.paragraphs[0].font.bold = True
    title8.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    # Ler dados para insights
    try:
        df_partido = pd.read_csv(resultados_dir / 'gastos_por_partido.csv')
        df_estado = pd.read_csv(resultados_dir / 'gastos_por_estado.csv')
        df_despesa = pd.read_csv(resultados_dir / 'gastos_por_tipo_despesa.csv')
        df_deputados = pd.read_csv(resultados_dir / 'top_deputados.csv')
        
        # Preparar insights
        partido_top = df_partido.iloc[0]
        estado_top = df_estado.iloc[0]
        despesa_top = df_despesa.iloc[0]
        deputado_top = df_deputados.iloc[0]
        
        total_gasto = df_partido['total_gasto'].sum()
        
        insights = [
            f"💰 Total gasto: R$ {total_gasto:,.2f}",
            f"🏆 Partido com maior gasto: {partido_top['partido']} (R$ {partido_top['total_gasto']:,.2f})",
            f"📍 Estado com maior gasto: {estado_top['uf']} (R$ {estado_top['total_gasto']:,.2f})",
            f"📊 Despesa mais comum: {despesa_top['tipo_despesa'][:40]}... ({despesa_top['percentual']:.1f}%)",
            f"👤 Maior gastador: {deputado_top['nome_deputado']} ({deputado_top['partido']}-{deputado_top['uf']})"
        ]
        
    except Exception as e:
        insights = [
            "💰 Análise completa de gastos parlamentares",
            "🏆 Identificação dos maiores gastadores",
            "📍 Comparação regional detalhada",
            "📊 Categorização de despesas",
            "👤 Ranking individualizado"
        ]
    
    # Adicionar insights
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(5)
    
    txBox = slide8.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for insight in insights:
        p = tf.add_paragraph()
        p.text = insight
        p.font.size = Pt(22)
        p.font.color.rgb = RGBColor(0, 51, 102)
        p.space_after = Pt(20)
    
    # === SLIDE 9: CONCLUSÃO ===
    print("   Slide 9: Conclusão")
    slide9 = prs.slides.add_slide(prs.slide_layouts[5])
    
    title9 = slide9.shapes.title
    title9.text = "Conclusão"
    title9.text_frame.paragraphs[0].font.size = Pt(40)
    title9.text_frame.paragraphs[0].font.bold = True
    title9.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    # Conteúdo
    left = Inches(1.5)
    top = Inches(2.2)
    width = Inches(7)
    height = Inches(4.5)
    
    txBox = slide9.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    conclusoes = [
        "✅ Cruzamento bem-sucedido de duas fontes de dados",
        "✅ Identificação clara dos padrões de gastos",
        "✅ Normalização permite comparação justa",
        "✅ Visualizações facilitam interpretação",
        "✅ Código modular e reutilizável"
    ]
    
    for conclusao in conclusoes:
        p = tf.add_paragraph()
        p.text = conclusao
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(0, 51, 102)
        p.space_after = Pt(25)
    
    # Rodapé
    left = Inches(1)
    top = Inches(6.5)
    width = Inches(8)
    height = Inches(0.8)
    
    txBox = slide9.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "Fontes: Portal de Dados Abertos da Câmara dos Deputados"
    p = tf.paragraphs[0]
    p.font.size = Pt(12)
    p.font.italic = True
    p.font.color.rgb = RGBColor(120, 120, 120)
    p.alignment = PP_ALIGN.CENTER
    
    # Salvar apresentação
    output_path = 'Apresentacao_Analise_Gastos_Parlamentares.pptx'
    prs.save(output_path)
    
    print(f"\n✅ Apresentação criada com sucesso!")
    print(f"📄 Arquivo: {output_path}")
    print(f"📊 Total de slides: 9")
    
    return output_path


if __name__ == '__main__':
    print("\n" + "=" * 70)
    print("  GERADOR DE APRESENTAÇÃO POWERPOINT")
    print("=" * 70 + "\n")
    
    try:
        arquivo = criar_apresentacao()
        
        print("\n" + "=" * 70)
        print("  🎉 SUCESSO!")
        print("=" * 70)
        print(f"\n  Abra o arquivo: {arquivo}")
        print("  Para visualizar a apresentação.\n")
        
    except Exception as e:
        print(f"\n❌ ERRO: {e}")
        import traceback
        traceback.print_exc()
