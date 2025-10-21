"""
Gera apresentação PowerPoint completa com análise de dados parlamentares.

Autor: Grupo 1 - Análise de Dados Governamentais
Data: Outubro 2025
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import pandas as pd
import os
from datetime import datetime


class ApresentacaoAnalise:
    """Classe para gerar apresentação PowerPoint completa."""
    
    def __init__(self, pasta_execucao, output_path=None):
        """
        Inicializa gerador de apresentação.
        
        Args:
            pasta_execucao: Caminho da pasta com resultados da execução
            output_path: Caminho customizado para salvar (opcional)
        """
        self.pasta_execucao = pasta_execucao
        self.output_path = output_path
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)
        
        # Cores do tema
        self.cor_titulo = RGBColor(0, 51, 102)  # Azul escuro
        self.cor_destaque = RGBColor(0, 102, 204)  # Azul médio
        self.cor_texto = RGBColor(51, 51, 51)  # Cinza escuro
        
    def criar_slide_titulo(self, titulo, subtitulo=""):
        """Cria slide de título."""
        slide_layout = self.prs.slide_layouts[6]  # Layout em branco
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Fundo azul
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0, 51, 102)
        
        # Título
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
        tf = txBox.text_frame
        tf.text = titulo
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        
        # Subtítulo
        if subtitulo:
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
            tf = txBox.text_frame
            tf.text = subtitulo
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(24)
            p.font.color.rgb = RGBColor(200, 200, 200)
        
        return slide
    
    def criar_slide_conteudo(self, titulo):
        """Cria slide de conteúdo padrão."""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Título
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        tf = txBox.text_frame
        tf.text = titulo
        p = tf.paragraphs[0]
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self.cor_titulo
        
        # Linha decorativa
        shape = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0.5), Inches(1.1),
            Inches(9), Inches(0.02)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.cor_destaque
        shape.line.color.rgb = self.cor_destaque
        
        return slide
    
    def adicionar_imagem(self, slide, caminho_imagem, left, top, width, height):
        """Adiciona imagem ao slide."""
        if os.path.exists(caminho_imagem):
            slide.shapes.add_picture(caminho_imagem, left, top, width=width, height=height)
            return True
        return False
    
    def adicionar_texto(self, slide, texto, left, top, width, height, tamanho=14, negrito=False, cor=None):
        """Adiciona caixa de texto ao slide."""
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.text = texto
        p = tf.paragraphs[0]
        p.font.size = Pt(tamanho)
        p.font.bold = negrito
        if cor:
            p.font.color.rgb = cor
        else:
            p.font.color.rgb = self.cor_texto
        return txBox
    
    def slide_1_capa(self):
        """Slide 1: Capa."""
        self.criar_slide_titulo(
            "Análise de Gastos Parlamentares",
            "Câmara dos Deputados - Cota Parlamentar"
        )
    
    def slide_2_equipe(self):
        """Slide 2: Equipe."""
        slide = self.criar_slide_conteudo("Equipe do Projeto")
        
        # Informações da equipe
        equipe = [
            ("Leticia", "21352"),
            ("Gabriel", "24734"),
            ("Thales", "24740"),
            ("Maria Fernanda", "24767")
        ]
        
        y_pos = 2.0
        for nome, matricula in equipe:
            self.adicionar_texto(
                slide, f"• {nome} - Matrícula {matricula}",
                Inches(2), Inches(y_pos), Inches(6), Inches(0.4),
                tamanho=20
            )
            y_pos += 0.6
        
        # Informações adicionais
        self.adicionar_texto(
            slide, "Grupo 1 - Ciência de Dados",
            Inches(2), Inches(5.5), Inches(6), Inches(0.4),
            tamanho=18, negrito=True, cor=self.cor_destaque
        )
        
        self.adicionar_texto(
            slide, "Outubro de 2025",
            Inches(2), Inches(6.0), Inches(6), Inches(0.4),
            tamanho=16, cor=RGBColor(100, 100, 100)
        )
    
    def slide_3_objetivos(self):
        """Slide 3: Objetivos."""
        slide = self.criar_slide_conteudo("Objetivos do Projeto")
        
        objetivos = [
            "Analisar gastos da Cota Parlamentar por partido político",
            "Comparar despesas entre estados brasileiros",
            "Identificar principais categorias de gastos",
            "Rankear deputados com maiores despesas",
            "Gerar insights sobre padrões de gastos públicos"
        ]
        
        y_pos = 2.0
        for obj in objetivos:
            self.adicionar_texto(
                slide, f"• {obj}",
                Inches(1), Inches(y_pos), Inches(8), Inches(0.5),
                tamanho=18
            )
            y_pos += 0.8
    
    def slide_4_metodologia(self):
        """Slide 4: Metodologia."""
        slide = self.criar_slide_conteudo("Metodologia")
        
        etapas = [
            ("1. Coleta de Dados", "CSV de despesas da Câmara dos Deputados"),
            ("2. Integração com API", "Dados cadastrais dos deputados"),
            ("3. Limpeza e Validação", "Remoção de registros inválidos"),
            ("4. Cruzamento de Dados", "Match por nome normalizado (>95%)"),
            ("5. Análise Estatística", "Agregações e rankings"),
            ("6. Visualização", "Gráficos profissionais em alta resolução")
        ]
        
        y_pos = 1.8
        for etapa, desc in etapas:
            self.adicionar_texto(
                slide, etapa,
                Inches(1), Inches(y_pos), Inches(8), Inches(0.3),
                tamanho=16, negrito=True, cor=self.cor_destaque
            )
            self.adicionar_texto(
                slide, desc,
                Inches(1.5), Inches(y_pos + 0.3), Inches(7.5), Inches(0.3),
                tamanho=14
            )
            y_pos += 0.8
    
    def slide_5_grafico_partidos(self):
        """Slide 5: Gráfico de Gastos por Partido."""
        slide = self.criar_slide_conteudo("Gastos por Partido Político")
        
        img_path = os.path.join(self.pasta_execucao, "gastos_por_partido.png")
        if self.adicionar_imagem(slide, img_path, Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5)):
            # Adicionar nota
            self.adicionar_texto(
                slide, "Valores em Reais (R$) - Top 10 partidos com maiores gastos",
                Inches(0.8), Inches(6.8), Inches(8.4), Inches(0.3),
                tamanho=12, cor=RGBColor(100, 100, 100)
            )
    
    def slide_6_grafico_estados(self):
        """Slide 6: Gráfico de Gastos por Estado."""
        slide = self.criar_slide_conteudo("Gastos por Estado (UF)")
        
        img_path = os.path.join(self.pasta_execucao, "gastos_por_estado.png")
        if self.adicionar_imagem(slide, img_path, Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5)):
            self.adicionar_texto(
                slide, "Valores em Reais (R$) - Top 10 estados com maiores gastos",
                Inches(0.8), Inches(6.8), Inches(8.4), Inches(0.3),
                tamanho=12, cor=RGBColor(100, 100, 100)
            )
    
    def slide_7_grafico_despesas(self):
        """Slide 7: Gráfico de Tipos de Despesa."""
        slide = self.criar_slide_conteudo("Principais Tipos de Despesa")
        
        img_path = os.path.join(self.pasta_execucao, "tipos_despesa.png")
        if self.adicionar_imagem(slide, img_path, Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5)):
            self.adicionar_texto(
                slide, "Valores em Reais (R$) - Top 15 categorias de despesas",
                Inches(0.8), Inches(6.8), Inches(8.4), Inches(0.3),
                tamanho=12, cor=RGBColor(100, 100, 100)
            )
    
    def slide_8_grafico_deputados(self):
        """Slide 8: Gráfico Top Deputados."""
        slide = self.criar_slide_conteudo("Top 20 Deputados - Maiores Gastos")
        
        img_path = os.path.join(self.pasta_execucao, "top_deputados.png")
        if self.adicionar_imagem(slide, img_path, Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5)):
            self.adicionar_texto(
                slide, "Ranking dos 20 deputados com maiores volumes de gastos individuais",
                Inches(0.8), Inches(6.8), Inches(8.4), Inches(0.3),
                tamanho=12, cor=RGBColor(100, 100, 100)
            )
    
    def slide_9_resumo_visual(self):
        """Slide 9: Dashboard Resumo."""
        slide = self.criar_slide_conteudo("Dashboard - Visão Geral")
        
        img_path = os.path.join(self.pasta_execucao, "resumo_geral.png")
        self.adicionar_imagem(slide, img_path, Inches(0.5), Inches(1.5), Inches(9), Inches(5.7))
    
    def slide_10_insights_partidos(self):
        """Slide 10: Insights sobre Partidos."""
        slide = self.criar_slide_conteudo("Insights - Gastos por Partido")
        
        # Tentar carregar dados
        csv_path = os.path.join(self.pasta_execucao, "gastos_por_partido.csv")
        if os.path.exists(csv_path):
            df = pd.read_csv(csv_path)
            
            if len(df) > 0:
                top1 = df.iloc[0]
                total_geral = df['valor_total'].sum()
                percentual_top1 = (top1['valor_total'] / total_geral * 100)
                
                insights = [
                    f"Partido com maior gasto: {top1['partido']}",
                    f"Valor total: R$ {top1['valor_total']:,.2f}",
                    f"Representa {percentual_top1:.1f}% do total geral",
                    f"Total de partidos analisados: {len(df)}",
                    f"Gasto médio por partido: R$ {df['valor_total'].mean():,.2f}"
                ]
                
                y_pos = 2.2
                for insight in insights:
                    self.adicionar_texto(
                        slide, f"• {insight}",
                        Inches(1.5), Inches(y_pos), Inches(7), Inches(0.4),
                        tamanho=18, cor=self.cor_texto
                    )
                    y_pos += 0.7
    
    def slide_11_insights_estados(self):
        """Slide 11: Insights sobre Estados."""
        slide = self.criar_slide_conteudo("Insights - Gastos por Estado")
        
        csv_path = os.path.join(self.pasta_execucao, "gastos_por_estado.csv")
        if os.path.exists(csv_path):
            df = pd.read_csv(csv_path)
            
            if len(df) > 0:
                top1 = df.iloc[0]
                total_geral = df['valor_total'].sum()
                percentual_top1 = (top1['valor_total'] / total_geral * 100)
                
                insights = [
                    f"Estado com maior gasto: {top1['estado']}",
                    f"Valor total: R$ {top1['valor_total']:,.2f}",
                    f"Representa {percentual_top1:.1f}% do total nacional",
                    f"Total de estados: {len(df)}",
                    f"Gasto médio por estado: R$ {df['valor_total'].mean():,.2f}"
                ]
                
                y_pos = 2.2
                for insight in insights:
                    self.adicionar_texto(
                        slide, f"• {insight}",
                        Inches(1.5), Inches(y_pos), Inches(7), Inches(0.4),
                        tamanho=18, cor=self.cor_texto
                    )
                    y_pos += 0.7
    
    def slide_12_insights_deputados(self):
        """Slide 12: Insights sobre Deputados."""
        slide = self.criar_slide_conteudo("Insights - Top Deputados")
        
        csv_path = os.path.join(self.pasta_execucao, "top_deputados.csv")
        if os.path.exists(csv_path):
            df = pd.read_csv(csv_path)
            
            if len(df) > 0:
                top1 = df.iloc[0]
                
                insights = [
                    f"Deputado com maior gasto: {top1['nome']}",
                    f"Partido: {top1['partido']} | Estado: {top1['estado']}",
                    f"Valor total: R$ {top1['valor_total']:,.2f}",
                    f"Média do Top 20: R$ {df['valor_total'].mean():,.2f}",
                    f"Amplitude: R$ {df['valor_total'].max() - df['valor_total'].min():,.2f}"
                ]
                
                y_pos = 2.2
                for insight in insights:
                    self.adicionar_texto(
                        slide, f"• {insight}",
                        Inches(1.5), Inches(y_pos), Inches(7), Inches(0.4),
                        tamanho=18, cor=self.cor_texto
                    )
                    y_pos += 0.7
    
    def slide_13_tabela_top5(self):
        """Slide 13: Tabela Top 5 Partidos."""
        slide = self.criar_slide_conteudo("Top 5 Partidos - Detalhamento")
        
        csv_path = os.path.join(self.pasta_execucao, "gastos_por_partido.csv")
        if os.path.exists(csv_path):
            df = pd.read_csv(csv_path).head(5)
            
            # Criar tabela
            rows, cols = len(df) + 1, 2
            left = Inches(2)
            top = Inches(2)
            width = Inches(6)
            height = Inches(3.5)
            
            table = slide.shapes.add_table(rows, cols, left, top, width, height).table
            
            # Cabeçalhos
            table.cell(0, 0).text = "Partido"
            table.cell(0, 1).text = "Valor Total (R$)"
            
            # Formatar cabeçalhos
            for col in range(cols):
                cell = table.cell(0, col)
                cell.fill.solid()
                cell.fill.fore_color.rgb = self.cor_destaque
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.bold = True
                paragraph.font.size = Pt(16)
                paragraph.font.color.rgb = RGBColor(255, 255, 255)
                paragraph.alignment = PP_ALIGN.CENTER
            
            # Dados
            for i, row in df.iterrows():
                table.cell(i + 1, 0).text = str(row['partido'])
                table.cell(i + 1, 1).text = f"R$ {row['valor_total']:,.2f}"
                
                # Formatar células
                for col in range(cols):
                    cell = table.cell(i + 1, col)
                    paragraph = cell.text_frame.paragraphs[0]
                    paragraph.font.size = Pt(14)
                    paragraph.alignment = PP_ALIGN.CENTER if col == 1 else PP_ALIGN.LEFT
    
    def slide_14_conclusoes(self):
        """Slide 14: Conclusões."""
        slide = self.criar_slide_conteudo("Conclusões")
        
        conclusoes = [
            "Identificação clara dos padrões de gastos parlamentares",
            "Diferenças significativas entre partidos e estados",
            "Concentração de gastos em categorias específicas",
            "Taxa de identificação de deputados superior a 95%",
            "Pipeline automatizado facilita análises futuras"
        ]
        
        y_pos = 2.2
        for conclusao in conclusoes:
            self.adicionar_texto(
                slide, f"✓ {conclusao}",
                Inches(1.5), Inches(y_pos), Inches(7), Inches(0.4),
                tamanho=18, cor=self.cor_texto
            )
            y_pos += 0.8
    
    def slide_15_encerramento(self):
        """Slide 15: Encerramento."""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Fundo azul
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0, 51, 102)
        
        # Título
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
        tf = txBox.text_frame
        tf.text = "Obrigado!"
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        
        # Subtítulo
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(0.8))
        tf = txBox.text_frame
        tf.text = "Grupo 1 - Análise de Dados Governamentais"
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(200, 200, 200)
        
        # Data
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(0.5))
        tf = txBox.text_frame
        tf.text = "Outubro de 2025"
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(150, 150, 150)
    
    def gerar(self):
        """Gera apresentação completa."""
        print("\n🎨 Gerando apresentação PowerPoint...")
        print("=" * 60)
        
        slides_info = [
            ("Slide 1/15", "Capa", self.slide_1_capa),
            ("Slide 2/15", "Equipe", self.slide_2_equipe),
            ("Slide 3/15", "Objetivos", self.slide_3_objetivos),
            ("Slide 4/15", "Metodologia", self.slide_4_metodologia),
            ("Slide 5/15", "Gráfico Partidos", self.slide_5_grafico_partidos),
            ("Slide 6/15", "Gráfico Estados", self.slide_6_grafico_estados),
            ("Slide 7/15", "Gráfico Despesas", self.slide_7_grafico_despesas),
            ("Slide 8/15", "Gráfico Deputados", self.slide_8_grafico_deputados),
            ("Slide 9/15", "Dashboard Resumo", self.slide_9_resumo_visual),
            ("Slide 10/15", "Insights Partidos", self.slide_10_insights_partidos),
            ("Slide 11/15", "Insights Estados", self.slide_11_insights_estados),
            ("Slide 12/15", "Insights Deputados", self.slide_12_insights_deputados),
            ("Slide 13/15", "Tabela Top 5", self.slide_13_tabela_top5),
            ("Slide 14/15", "Conclusões", self.slide_14_conclusoes),
            ("Slide 15/15", "Encerramento", self.slide_15_encerramento)
        ]
        
        for num, nome, metodo in slides_info:
            print(f"  {num}: {nome}...")
            metodo()
        
        # Salvar
        if self.output_path:
            output_file = self.output_path
        else:
            output_file = os.path.join(self.pasta_execucao, "Apresentacao_Completa.pptx")
        
        self.prs.save(output_file)
        
        print("=" * 60)
        print(f"✅ Apresentação salva em: {output_file}")
        print(f"📊 Total de slides: 15")
        print(f"📁 Tamanho: {os.path.getsize(output_file) / 1024 / 1024:.2f} MB\n")
        
        return output_file


def main():
    """Função principal para execução standalone."""
    # Encontrar pasta de execução mais recente
    resultados_dir = os.path.join(os.path.dirname(__file__), "resultados")
    
    if not os.path.exists(resultados_dir):
        print("❌ Pasta 'resultados' não encontrada!")
        print("💡 Execute primeiro: python src/main.py dados/Ano-2025.csv")
        return
    
    # Listar pastas de execução
    execucoes = [d for d in os.listdir(resultados_dir) if d.startswith("execucao_")]
    
    if not execucoes:
        print("❌ Nenhuma execução encontrada!")
        print("💡 Execute primeiro: python src/main.py dados/Ano-2025.csv")
        return
    
    # Pegar mais recente
    execucoes.sort(reverse=True)
    pasta_execucao = os.path.join(resultados_dir, execucoes[0])
    
    print(f"\n📂 Usando resultados de: {execucoes[0]}")
    
    # Gerar apresentação
    apresentacao = ApresentacaoAnalise(pasta_execucao)
    apresentacao.gerar()
    
    print("✅ Apresentação gerada com sucesso!")
    print(f"📂 Localização: {pasta_execucao}")


if __name__ == "__main__":
    main()
